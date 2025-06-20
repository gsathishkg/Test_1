
import streamlit as st
from graphviz import Digraph
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import json
import os
import tempfile

st.set_page_config(page_title="Flowchart Designer", layout="wide")

if "steps" not in st.session_state:
    st.session_state.steps = []
if "edit_index" not in st.session_state:
    st.session_state.edit_index = None

st.title("🎨 Flowchart Designer")

def generate_flowchart(steps):
    if not steps:
        return None

    dot = Digraph(format="png")
    dot.attr(rankdir='TB', splines='ortho', nodesep='1', ranksep='1')
    step_map = {s["id"]: s for s in steps}
    added = set()

    def add_node(nid):
        if nid in added:
            return
        step = step_map.get(nid, {"label": nid, "type": "process"})
        shape = {
            "start": "circle",
            "end": "doublecircle",
            "decision": "diamond",
            "process": "box"
        }.get(step.get("type", "process"), "box")

        fillcolor = {
            "start": "#007acc",
            "end": "#2e7d32",
            "decision": "#fdd835",
            "process": "#90a4ae"
        }.get(step.get("type", "process"), "#e0e0e0")

        style = "rounded,filled,bold"

        dot.node(
            nid,
            f"<<b>{step.get('label', nid)}</b>>",
            shape=shape,
            style=style,
            fillcolor=fillcolor,
            fontcolor="black",
            fontsize="12",
            fontname="Verdana",
        )
        added.add(nid)

    for step in steps:
        sid = step.get("id")
        if not sid:
            continue
        add_node(sid)
        for nxt in step.get("next", []):
            tid = nxt.get("id")
            if tid:
                add_node(tid)
                label = nxt.get("label", "")
                color = {
                    "yes": "green",
                    "no": "red"
                }.get(label.lower(), "black")
                dot.edge(
                    sid,
                    tid,
                    label=label,
                    fontsize="10",
                    fontname="Verdana",
                    color=color
                )

    return dot if added else None

# Input Form
with st.form("step_form", clear_on_submit=True):
    if st.session_state.edit_index is None:
        st.subheader("➕ Add New Step")
        default = {"id": "", "label": "", "type": "process", "next": []}
    else:
        st.subheader("✏️ Edit Step")
        default = st.session_state.steps[st.session_state.edit_index]

    step_id = st.text_input("Step ID", value=default["id"])
    label = st.text_input("Label", value=default["label"])
    step_type = st.selectbox("Step Type", ["start", "process", "decision", "end"],
                             index=["start", "process", "decision", "end"].index(default["type"]))
    next_steps = []
    for i in range(2):
        col1, col2 = st.columns(2)
        with col1:
            nid = st.text_input(f"Next Step ID {i+1}", value=default["next"][i]["id"] if i < len(default["next"]) else "")
        with col2:
            lbl = st.text_input(f"Label (e.g., Yes/No) {i+1}", value=default["next"][i].get("label", "") if i < len(default["next"]) else "")
        if nid:
            next_steps.append({"id": nid, "label": lbl})

    submitted = st.form_submit_button("✅ Save Step")
    if submitted:
        new_step = {"id": step_id, "label": label, "type": step_type, "next": next_steps}
        if st.session_state.edit_index is None:
            st.session_state.steps.append(new_step)
        else:
            st.session_state.steps[st.session_state.edit_index] = new_step
            st.session_state.edit_index = None
        st.success("Step saved!")

# Edit/Delete UI
st.subheader("🧾 Flow Steps")
for idx, step in enumerate(st.session_state.steps):
    cols = st.columns([4, 1, 1])
    with cols[0]:
        st.code(json.dumps(step, indent=2))
    with cols[1]:
        if st.button("✏️ Edit", key=f"edit_{idx}"):
            st.session_state.edit_index = idx
    with cols[2]:
        if st.button("🗑 Delete", key=f"delete_{idx}"):
            del st.session_state.steps[idx]
            if st.session_state.edit_index == idx:
                st.session_state.edit_index = None
            st.experimental_rerun()

# Preview
st.subheader("🖼 Flowchart Preview")
flow = generate_flowchart(st.session_state.steps)

if flow:
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmpfile:
            flow.render(tmpfile.name, format="png", cleanup=True)
            img = Image.open(tmpfile.name + ".png")
            st.image(img, caption="Generated Flowchart", use_column_width=True)
    except Exception as e:
        st.error(f"❌ Failed to render image: {e}")
else:
    st.warning("No steps defined yet.")

# Export
col1, col2, col3 = st.columns(3)

with col1:
    if flow:
        st.download_button("⬇️ Download PNG", data=flow.pipe(format="png"),
                           file_name="flowchart.png", mime="image/png")

with col2:
    if flow and st.button("📊 Export to PPTX"):
        flow.render(filename="temp_flow", format="png", cleanup=True)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture("temp_flow.png", Inches(1), Inches(1), height=Inches(5.5))
        prs.save("flowchart_output.pptx")
        with open("flowchart_output.pptx", "rb") as f:
            st.download_button("📥 Download PPTX", data=f, file_name="flowchart.pptx")

with col3:
    with st.expander("💾 Save / Load"):
        st.download_button("📥 Download JSON", json.dumps(st.session_state.steps, indent=2),
                           file_name="flow.json")
        uploaded = st.file_uploader("Upload flow JSON", type=["json"])
        if uploaded:
            st.session_state.steps = json.load(uploaded)
            st.experimental_rerun()
